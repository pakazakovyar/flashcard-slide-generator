from http.client import responses
from typing import List
from fastapi import FastAPI, staticfiles, Request, Form, File, UploadFile
from fastapi.templating import Jinja2Templates
from fastapi.responses import HTMLResponse, RedirectResponse, FileResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image as PILImage
import io

PIXELS_WIDTH = 1920
PIXELS_HEIGHT = 1080
DPI = 96

EMU_PER_INCH = 914400
WIDTH_EMU = int((PIXELS_WIDTH / DPI) * EMU_PER_INCH)
HEIGHT_EMU = int((PIXELS_HEIGHT / DPI) * EMU_PER_INCH)


def create_presentation_with_images_and_words(image_binaries, words, output_path='presentation.pptx'):
    """
    Создает презентацию с изображениями из бинарных строк и словами

    Args:
        image_binaries: список бинарных строк (bytes) изображений
        words: список слов для каждого изображения
        output_path: путь для сохранения презентации
    """
    prs = Presentation()
    prs.slide_width = WIDTH_EMU
    prs.slide_height = HEIGHT_EMU
    blank_slide_layout = prs.slide_layouts[6]

    for img_binary, word in zip(image_binaries, words):
        with PILImage.open(io.BytesIO(img_binary)) as img:
            img_width_px, img_height_px = img.size
            img_dpi = img.info.get('dpi', (DPI, DPI))[0] or DPI

        img_width_emu = int((img_width_px / img_dpi) * EMU_PER_INCH)
        img_height_emu = int((img_height_px / img_dpi) * EMU_PER_INCH)

        scale = min(WIDTH_EMU / img_width_emu, HEIGHT_EMU / img_height_emu)
        new_width = int(img_width_emu * scale)
        new_height = int(img_height_emu * scale)

        left_img = (WIDTH_EMU - new_width) // 2
        top_img = (HEIGHT_EMU - new_height) // 2

        slide1 = prs.slides.add_slide(blank_slide_layout)
        slide1.shapes.add_picture(
            io.BytesIO(img_binary),
            left_img,
            top_img,
            width=new_width,
            height=new_height
        )

        slide2 = prs.slides.add_slide(blank_slide_layout)
        slide2.shapes.add_picture(
            io.BytesIO(img_binary),
            left_img,
            top_img,
            width=new_width,
            height=new_height
        )

        text_left = 0
        text_width = WIDTH_EMU
        text_height = Inches(1.2)
        text_top = (HEIGHT_EMU - text_height) // 2

        def add_text(slide, text, color, offset_x=0, offset_y=0):
            box = slide.shapes.add_textbox(
                left=text_left + offset_x,
                top=text_top + offset_y,
                width=text_width,
                height=text_height
            )
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(60)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

        offsets = [(dx, dy) for dx in (-2, 0, 2) for dy in (-2, 0, 2) if not (dx == 0 and dy == 0)]
        for dx, dy in offsets:
            add_text(slide2, word, RGBColor(0, 0, 0), Inches(dx / 96), Inches(dy / 96))

        add_text(slide2, word, RGBColor(255, 255, 255))

    prs.save(output_path)










app = FastAPI()
templates = Jinja2Templates(directory="./../templates")

app.words = []
app.images = []

@app.get("/")
def get_main_page(request: Request):
    return templates.TemplateResponse(
        "index.html",
        {"request": request}
    )
@app.post("/",)
def handle_words(
        request: Request,
        text: str = Form(...)
):
    app.words = text.split(";")
    print(app.words)
    return RedirectResponse(url="/images", status_code=303)



@app.get("/images")
def set_images(request: Request):
    return templates.TemplateResponse(
        "words.html",
        {"request": request,
         "words" : app.words}
    )


@app.post("/upload-images/")
async def upload_images(request: Request, words: List[str] = Form(...), images: List[UploadFile] = File(...)):
    app.images = [await image.read() for image in images]
    return RedirectResponse(url="/presentation", status_code=303)


@app.get("/presentation")
def get_presentation(request: Request):
    create_presentation_with_images_and_words(app.images, app.words)
    return FileResponse(
        path="presentation.pptx",
        filename="presentation.pptx",
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
