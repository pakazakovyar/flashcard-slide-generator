from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

PIXELS_WIDTH = 1920
PIXELS_HEIGHT = 1080
DPI = 96

EMU_PER_INCH = 914400
WIDTH_EMU = int((PIXELS_WIDTH / DPI) * EMU_PER_INCH)
HEIGHT_EMU = int((PIXELS_HEIGHT / DPI) * EMU_PER_INCH)

def create_presentation_with_images_and_words(image_paths, words, output_path='presentation.pptx'):

    prs = Presentation()
    prs.slide_width = WIDTH_EMU
    prs.slide_height = HEIGHT_EMU
    blank_slide_layout = prs.slide_layouts[6]  # Пустой слайд

    for img_path, word in zip(image_paths, words):
        with PILImage.open(img_path) as img:
            img_width_px, img_height_px = img.size
            img_dpi = img.info.get('dpi', (DPI, DPI))[0] or DPI

        img_width_emu = int((img_width_px / img_dpi) * EMU_PER_INCH)
        img_height_emu = int((img_height_px / img_dpi) * EMU_PER_INCH)

        scale = min(WIDTH_EMU / img_width_emu, HEIGHT_EMU / img_height_emu)
        new_width = int(img_width_emu * scale)
        new_height = int(img_height_emu * scale)

        left_img = (WIDTH_EMU - new_width) // 2
        top_img = (HEIGHT_EMU - new_height) // 2

        slide1 = prs.slides.add_slide(blank_slide_layout)
        slide1.shapes.add_picture(img_path, left_img, top_img, width=new_width, height=new_height)

        slide2 = prs.slides.add_slide(blank_slide_layout)
        slide2.shapes.add_picture(img_path, left_img, top_img, width=new_width, height=new_height)

        text_left = 0
        text_width = WIDTH_EMU
        text_height = Inches(1.2)
        text_top = (HEIGHT_EMU - text_height) // 2  # по вертикали

        def add_text(slide, text, color, offset_x=0, offset_y=0):
            box = slide.shapes.add_textbox(
                left=text_left + offset_x,
                top=text_top + offset_y,
                width=text_width,
                height=text_height
            )
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            p.text = text
            p.font.size = Pt(60)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER

        offsets = [(dx, dy) for dx in (-2, 0, 2) for dy in (-2, 0, 2) if not (dx == 0 and dy == 0)]
        for dx, dy in offsets:
            add_text(slide2, word, RGBColor(0, 0, 0), Inches(dx / 96), Inches(dy / 96))

        add_text(slide2, word, RGBColor(255, 255, 255))

    prs.save(output_path)

images = ['img_1.png', 'img.png', "img_2.png"]
words = ['Солнце', 'Дождь', "asfasf"]
create_presentation_with_images_and_words(images, words, 'my_presentation2.pptx')